"""Generador de materiales de capacitación TherapyCord (24 abr 2026).

Produce:
- docs/capacitacion/TherapyCord_Capacitacion_Terapeutas.pptx (presentación ~50 min)
- docs/capacitacion/TherapyCord_Cheatsheet_Terapeutas.pdf  (handout 1 hoja, doble cara)
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from reportlab.lib.pagesizes import LETTER
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor
from reportlab.pdfgen import canvas

OUT_DIR = Path(__file__).resolve().parent.parent / "docs" / "capacitacion"
OUT_DIR.mkdir(parents=True, exist_ok=True)

TEAL = RGBColor(0x0E, 0x7C, 0x86)
TEAL_DARK = RGBColor(0x08, 0x4C, 0x53)
CORAL = RGBColor(0xE8, 0x6A, 0x4C)
GRAY_DARK = RGBColor(0x2B, 0x2B, 0x2B)
GRAY_MID = RGBColor(0x55, 0x55, 0x55)
GRAY_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_background(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=GRAY_DARK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_rule(slide, left, top, width, color=TEAL, height_pt=3):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(height_pt)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def add_card(slide, left, top, width, height, fill=GRAY_LIGHT, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    return shape


def add_header(slide, title, kicker=None):
    if kicker:
        add_text(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.3),
                 kicker.upper(), size=12, bold=True, color=TEAL)
    add_text(slide, Inches(0.6), Inches(0.6), Inches(12), Inches(0.9),
             title, size=32, bold=True, color=TEAL_DARK)
    add_rule(slide, Inches(0.6), Inches(1.55), Inches(1.2))


# ===================== PPTX =====================

def build_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ----- Slide 1 — Portada
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    # Banda lateral izquierda teal
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.6), Inches(7.5))
    band.line.fill.background()
    band.fill.solid(); band.fill.fore_color.rgb = TEAL
    add_text(s, Inches(1.0), Inches(2.2), Inches(11), Inches(0.5),
             "CAPACITACIÓN INTERNA", size=14, bold=True, color=TEAL)
    add_text(s, Inches(1.0), Inches(2.7), Inches(11), Inches(1.5),
             "TherapyCord: cómo funciona\ntu sistema de agenda",
             size=44, bold=True, color=TEAL_DARK)
    add_text(s, Inches(1.0), Inches(4.6), Inches(11), Inches(0.5),
             "Salud Total + Cordelia (voz y WhatsApp)",
             size=20, color=GRAY_MID)
    add_rule(s, Inches(1.0), Inches(5.3), Inches(1.5), color=CORAL)
    add_text(s, Inches(1.0), Inches(5.5), Inches(11), Inches(0.4),
             "24 de abril 2026 · Hospital Ángeles México",
             size=14, color=GRAY_MID)
    add_text(s, Inches(1.0), Inches(5.9), Inches(11), Inches(0.4),
             "Dirigido a: Lic. Lilia · Lic. Montserrat · Lic. Kevin · Lic. Harold",
             size=14, color=GRAY_MID)

    # ----- Slide 2 — Agenda del día
    s = prs.slides.add_slide(blank); set_background(s, WHITE)
    add_header(s, "Lo que veremos hoy (≈50 min)", kicker="Agenda")
    items = [
        ("1", "El sistema en 1 minuto", "5 min"),
        ("2", "Cómo se agenda un paciente ahora", "10 min"),
        ("3", "Tu rol: antes, durante y después de la cita", "10 min"),
        ("4", "Salud Total: acceso y tu agenda personal", "15 min"),
        ("5", "Qué hacer si algo falla", "5 min"),
        ("6", "Preguntas y respuestas", "5 min"),
    ]
    top = Inches(2.0)
    for i, (n, title, dur) in enumerate(items):
        y = top + Inches(i * 0.75)
        add_card(s, Inches(0.8), y, Inches(0.7), Inches(0.6), fill=TEAL)
        add_text(s, Inches(0.8), y, Inches(0.7), Inches(0.6),
                 n, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.7), y + Inches(0.08), Inches(8.5), Inches(0.5),
                 title, size=20, bold=True, color=GRAY_DARK)
        add_text(s, Inches(11.0), y + Inches(0.1), Inches(1.8), Inches(0.5),
                 dur, size=14, color=CORAL, bold=True, align=PP_ALIGN.RIGHT)

    # ----- Slide 3 — El sistema en 1 minuto
    s = prs.slides.add_slide(blank); set_background(s, WHITE)
    add_header(s, "El sistema en un minuto", kicker="Módulo 1")
    add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.6),
             "TherapyCord opera con DOS piezas que trabajan juntas:",
             size=18, color=GRAY_MID)

    # Two cards
    card_w, card_h = Inches(5.8), Inches(4.2)
    add_card(s, Inches(0.6), Inches(2.6), card_w, card_h, fill=GRAY_LIGHT)
    add_text(s, Inches(1.0), Inches(2.85), Inches(5.0), Inches(0.5),
             "1. SALUD TOTAL", size=14, bold=True, color=TEAL)
    add_text(s, Inches(1.0), Inches(3.25), Inches(5.0), Inches(0.8),
             "Expedientes y agenda médica", size=22, bold=True, color=GRAY_DARK)
    for i, line in enumerate([
        "• Tu agenda personal con citas del día",
        "• Expediente de cada paciente",
        "• Notas de sesión",
        "• Cada terapeuta tiene su propia cuenta",
    ]):
        add_text(s, Inches(1.0), Inches(4.3 + i*0.4), Inches(5.0), Inches(0.4),
                 line, size=15, color=GRAY_DARK)

    add_card(s, Inches(6.9), Inches(2.6), card_w, card_h, fill=TEAL)
    add_text(s, Inches(7.3), Inches(2.85), Inches(5.0), Inches(0.5),
             "2. CORDELIA", size=14, bold=True, color=WHITE)
    add_text(s, Inches(7.3), Inches(3.25), Inches(5.0), Inches(0.8),
             "Asistente virtual", size=22, bold=True, color=WHITE)
    for i, line in enumerate([
        "• Contesta el teléfono automáticamente",
        "• Atiende mensajes de WhatsApp",
        "• Agenda citas en Salud Total por ti",
        "• Sólo agenda a PACIENTES EXISTENTES",
    ]):
        add_text(s, Inches(7.3), Inches(4.3 + i*0.4), Inches(5.0), Inches(0.4),
                 line, size=15, color=WHITE)

    # ----- Slide 4 — Dos canales
    s = prs.slides.add_slide(blank); set_background(s, WHITE)
    add_header(s, "Cordelia tiene dos canales", kicker="Cómo te contactan los pacientes")

    # Phone card
    add_card(s, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4.4), fill=GRAY_LIGHT)
    add_text(s, Inches(1.1), Inches(2.5), Inches(5.2), Inches(0.8),
             "📞  POR TELÉFONO", size=22, bold=True, color=TEAL_DARK)
    add_text(s, Inches(1.1), Inches(3.3), Inches(5.2), Inches(0.4),
             "55 6304 9089 (línea clínica)", size=16, bold=True, color=GRAY_DARK)
    add_text(s, Inches(1.1), Inches(3.8), Inches(5.2), Inches(2.5),
             "• Cordelia contesta con su voz\n"
             "• Agenda, reagenda y cancela\n"
             "• Responde dudas frecuentes\n"
             "• Transfiere a humano si hace falta",
             size=15, color=GRAY_DARK)

    add_card(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.4), fill=GRAY_LIGHT)
    add_text(s, Inches(7.1), Inches(2.5), Inches(5.2), Inches(0.8),
             "💬  POR WHATSAPP", size=22, bold=True, color=TEAL_DARK)
    add_text(s, Inches(7.1), Inches(3.3), Inches(5.2), Inches(0.4),
             "55 2884 1932 (WhatsApp clínica)", size=16, bold=True, color=GRAY_DARK)
    add_text(s, Inches(7.1), Inches(3.8), Inches(5.2), Inches(2.5),
             "• Cordelia responde por escrito\n"
             "• Mismo flujo que por llamada\n"
             "• Envía código de confirmación\n"
             "• Manda recordatorio un día antes",
             size=15, color=GRAY_DARK)

    add_text(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
             "Los dos canales usan el MISMO expediente en Salud Total.",
             size=14, bold=True, color=CORAL, align=PP_ALIGN.CENTER)

    # ----- Slide 5 — Los dos números clave
    s = prs.slides.add_slide(blank); set_background(s, WHITE)
    add_header(s, "Dos números que tienes que recordar", kicker="Teléfonos clave")

    add_card(s, Inches(0.8), Inches(2.1), Inches(5.8), Inches(4.6), fill=TEAL)
    add_text(s, Inches(1.1), Inches(2.4), Inches(5.2), Inches(0.4),
             "ADMISIÓN", size=14, bold=True, color=WHITE)
    add_text(s, Inches(1.1), Inches(2.85), Inches(5.2), Inches(0.9),
             "55 5516 9900", size=36, bold=True, color=WHITE)
    add_text(s, Inches(1.1), Inches(3.8), Inches(5.2), Inches(0.5),
             "Para quién:", size=14, bold=True, color=WHITE)
    add_text(s, Inches(1.1), Inches(4.2), Inches(5.2), Inches(2.4),
             "• Paciente NUEVO (primera vez)\n"
             "• Paciente con aseguradora / GMM\n"
             "• Paciente con pase médico\n"
             "• Paciente que busca reembolso\n"
             "• Expediente que no se pudo validar",
             size=15, color=WHITE)

    add_card(s, Inches(6.8), Inches(2.1), Inches(5.8), Inches(4.6),
             fill=GRAY_LIGHT, line=CORAL)
    add_text(s, Inches(7.1), Inches(2.4), Inches(5.2), Inches(0.4),
             "CLÍNICA (celular)", size=14, bold=True, color=CORAL)
    add_text(s, Inches(7.1), Inches(2.85), Inches(5.2), Inches(0.9),
             "55 2884 1932", size=36, bold=True, color=TEAL_DARK)
    add_text(s, Inches(7.1), Inches(3.8), Inches(5.2), Inches(0.5),
             "Para quién:", size=14, bold=True, color=GRAY_DARK)
    add_text(s, Inches(7.1), Inches(4.2), Inches(5.2), Inches(2.4),
             "• Queja, molestia o frustración\n"
             "• Emergencia médica real\n"
             "• Transferencia a humano\n"
             "• Lo que Cordelia no puede resolver",
             size=15, color=GRAY_DARK)

    # ----- Slide 6 — ¿Quién agenda a quién?
    s = prs.slides.add_slide(blank); set_background(s, WHITE)
    add_header(s, "¿Quién agenda a quién?", kicker="Regla de oro")

    add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.5),
             "Cordelia sólo agenda a pacientes EXISTENTES ya validados.",
             size=18, bold=True, color=GRAY_DARK)

    # Existing flow
    add_card(s, Inches(0.6), Inches(2.6), Inches(6.0), Inches(4.0), fill=TEAL)
    add_text(s, Inches(0.9), Inches(2.85), Inches(5.4), Inches(0.5),
             "✅ PACIENTE EXISTENTE", size=16, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.4), Inches(5.4), Inches(0.6),
             "Cordelia lo agenda directo", size=22, bold=True, color=WHITE)
    for i, line in enumerate([
        "1. Pide teléfono",
        "2. Pide fecha de nacimiento O email",
        "3. Confirma en Salud Total",
        "4. Ofrece horarios con tu agenda",
        "5. Agenda y envía código MMDD-##",
    ]):
        add_text(s, Inches(0.9), Inches(4.2 + i*0.42), Inches(5.4), Inches(0.4),
                 line, size=14, color=WHITE)

    # New flow
    add_card(s, Inches(6.8), Inches(2.6), Inches(6.0), Inches(4.0), fill=CORAL)
    add_text(s, Inches(7.1), Inches(2.85), Inches(5.4), Inches(0.5),
             "🔁 PACIENTE NUEVO · ASEGURADORA · PASE · REEMBOLSO",
             size=12, bold=True, color=WHITE)
    add_text(s, Inches(7.1), Inches(3.4), Inches(5.4), Inches(0.6),
             "Cordelia deriva a ADMISIÓN", size=22, bold=True, color=WHITE)
    for i, line in enumerate([
        "1. Cordelia pregunta: «¿primera vez?»",
        "2. Si sí → da el número 55 5516 9900",
        "3. Ofrece conectar la llamada",
        "4. Admisión crea el expediente",
        "5. Admisión agenda la valoración inicial",
    ]):
        add_text(s, Inches(7.1), Inches(4.2 + i*0.42), Inches(5.4), Inches(0.4),
                 line, size=14, color=WHITE)

    # ----- Slide 7 — Validación del paciente existente
    s = prs.slides.add_slide(blank); set_background(s, WHITE)
    add_header(s, "Cómo Cordelia valida a un paciente existente", kicker="Seguridad del expediente")

    add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.5),
             "Antes de agendar, Cordelia pide DOS datos que coincidan con el expediente:",
             size=17, color=GRAY_DARK)

    steps = [
        ("1", "Teléfono", "Siempre obligatorio"),
        ("+", "Fecha de nacimiento", "O"),
        ("=", "Correo electrónico", "Cualquiera de los dos"),
    ]
    for i, (mark, title, sub) in enumerate(steps):
        x = Inches(0.6 + i*4.25)
        add_card(s, x, Inches(2.7), Inches(4.0), Inches(3.4), fill=GRAY_LIGHT)
        color = TEAL if mark != "+" else CORAL
        add_card(s, x + Inches(1.5), Inches(2.9), Inches(1.0), Inches(1.0), fill=color)
        add_text(s, x + Inches(1.5), Inches(2.9), Inches(1.0), Inches(1.0),
                 mark, size=34, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.3), Inches(4.1), Inches(3.4), Inches(0.6),
                 title, size=20, bold=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.3), Inches(4.8), Inches(3.4), Inches(0.5),
                 sub, size=14, color=GRAY_MID, align=PP_ALIGN.CENTER)

    add_card(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.7),
             fill=RGBColor(0xFF, 0xF5, 0xE9))
    add_text(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.7),
             "Si los datos no coinciden o el paciente no existe → Cordelia lo deriva a admisión.",
             size=14, bold=True, color=CORAL,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ----- Slide 8 — Tu rol como terapeuta
    s = prs.slides.add_slide(blank); set_background(s, WHITE)
    add_header(s, "Tu rol como terapeuta", kicker="Antes · Durante · Después")

    cols = [
        ("ANTES de la cita",
         ["Abre tu agenda en Salud Total",
          "Revisa nombre, hora y motivo",
          "Lee las notas de la sesión previa",
          "Prepara el material que necesites"]),
        ("DURANTE la sesión",
         ["40-45 min individualizados",
          "Pregunta por dolor o cambios",
          "Aplica el plan de tratamiento",
          "Explica ejercicios en casa"]),
        ("DESPUÉS de la sesión",
         ["Registra la nota en Salud Total",
          "Anota ejercicios realizados",
          "Plan para siguiente sesión",
          "Observaciones relevantes"]),
    ]
    for i, (title, lines) in enumerate(cols):
        x = Inches(0.6 + i*4.25)
        add_card(s, x, Inches(2.1), Inches(4.0), Inches(4.8), fill=GRAY_LIGHT)
        add_rule(s, x + Inches(0.3), Inches(2.35), Inches(3.4),
                 color=[TEAL, TEAL_DARK, CORAL][i], height_pt=4)
        add_text(s, x + Inches(0.3), Inches(2.55), Inches(3.4), Inches(0.6),
                 title, size=18, bold=True, color=GRAY_DARK)
        for j, line in enumerate(lines):
            add_text(s, x + Inches(0.3), Inches(3.3 + j*0.7), Inches(3.4), Inches(0.6),
                     "• " + line, size=14, color=GRAY_DARK)

    # ----- Slide 9 — Acceso a Salud Total
    s = prs.slides.add_slide(blank); set_background(s, WHITE)
    add_header(s, "Cómo entras a Salud Total", kicker="Módulo 4")

    add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.5),
             "Cada terapeuta tiene su propia cuenta. Úsala siempre, no la prestes.",
             size=16, color=GRAY_MID)

    add_card(s, Inches(0.6), Inches(2.6), Inches(12.1), Inches(1.1), fill=TEAL)
    add_text(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.1),
             "🌐   www.saludtotal.mx", size=26, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

    # Users table
    rows = [
        ("Lic. Lilia Salazar",   "lilia@therapycord.com",   "123456789"),
        ("Lic. Montserrat",      "monste@therapycord.com",  "123456789"),
        ("Lic. Kevin Castellanos","kevin@therapycord.com",  "123456789"),
        ("Lic. Harold Ildefonso","harold@therapycord.com",  "123456789"),
    ]
    header_y = Inches(4.0)
    add_card(s, Inches(0.6), header_y, Inches(12.1), Inches(0.5), fill=GRAY_DARK)
    add_text(s, Inches(0.9), header_y, Inches(4.0), Inches(0.5),
             "TERAPEUTA", size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.0), header_y, Inches(5.0), Inches(0.5),
             "USUARIO (EMAIL)", size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(10.2), header_y, Inches(2.4), Inches(0.5),
             "CONTRASEÑA", size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    for i, (name, email, pwd) in enumerate(rows):
        y = header_y + Inches(0.55 + i*0.5)
        fill = GRAY_LIGHT if i % 2 == 0 else WHITE
        add_card(s, Inches(0.6), y, Inches(12.1), Inches(0.45), fill=fill)
        add_text(s, Inches(0.9), y, Inches(4.0), Inches(0.45),
                 name, size=13, bold=True, color=GRAY_DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(5.0), y, Inches(5.0), Inches(0.45),
                 email, size=13, color=GRAY_DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(10.2), y, Inches(2.4), Inches(0.45),
                 pwd, size=13, color=CORAL, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4),
             "⚠ Cambia tu contraseña en cuanto entres por primera vez.",
             size=13, bold=True, color=CORAL, align=PP_ALIGN.CENTER)

    # ----- Slide 10 — Tu agenda del día
    s = prs.slides.add_slide(blank); set_background(s, WHITE)
    add_header(s, "Así ves tu agenda del día", kicker="Paso a paso")

    steps = [
        ("1", "Entra a Salud Total con tu usuario"),
        ("2", "Ve al menú 'Agenda'"),
        ("3", "Filtra por tu nombre (sólo tú)"),
        ("4", "Elige el día en el calendario"),
        ("5", "Verde = libre · Rojo = ocupado"),
        ("6", "Click en la cita → verás paciente y motivo"),
    ]
    for i, (n, txt) in enumerate(steps):
        row = i // 2; col = i % 2
        x = Inches(0.8 + col*6.3)
        y = Inches(2.3 + row*1.35)
        add_card(s, x, y, Inches(5.8), Inches(1.1), fill=GRAY_LIGHT)
        add_card(s, x + Inches(0.2), y + Inches(0.2), Inches(0.7), Inches(0.7), fill=TEAL)
        add_text(s, x + Inches(0.2), y + Inches(0.2), Inches(0.7), Inches(0.7),
                 n, size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(1.1), y, Inches(4.6), Inches(1.1),
                 txt, size=15, color=GRAY_DARK, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.5),
             "Tip: guarda el link directo a tu agenda como favorito en el navegador.",
             size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # ----- Slide 11 — Qué NO tienes que hacer
    s = prs.slides.add_slide(blank); set_background(s, WHITE)
    add_header(s, "Lo que YA NO es tu chamba", kicker="Cordelia lo hace por ti")

    items = [
        ("❌ Contestar el teléfono",
         "Cordelia atiende las llamadas entrantes."),
        ("❌ Responder mensajes de WhatsApp",
         "Cordelia platica con el paciente y agenda."),
        ("❌ Agendar pacientes nuevos",
         "Los nuevos van a admisión 55 5516 9900."),
        ("❌ Dar precios por teléfono",
         "Precios se explican en la valoración inicial."),
    ]
    for i, (title, desc) in enumerate(items):
        row = i // 2; col = i % 2
        x = Inches(0.8 + col*6.3)
        y = Inches(2.2 + row*2.2)
        add_card(s, x, y, Inches(5.8), Inches(1.9), fill=GRAY_LIGHT)
        add_text(s, x + Inches(0.4), y + Inches(0.3), Inches(5.2), Inches(0.6),
                 title, size=20, bold=True, color=CORAL)
        add_text(s, x + Inches(0.4), y + Inches(1.0), Inches(5.2), Inches(0.9),
                 desc, size=14, color=GRAY_DARK)

    add_text(s, Inches(0.6), Inches(6.8), Inches(12.1), Inches(0.4),
             "Tu atención se enfoca en lo clínico. La logística la resuelve el sistema.",
             size=14, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

    # ----- Slide 12 — Si algo falla
    s = prs.slides.add_slide(blank); set_background(s, WHITE)
    add_header(s, "Si algo falla", kicker="Escala inteligente")

    rows = [
        ("Cordelia no contesta una llamada", "Avisa a Vicci",         "55 1188 0301"),
        ("Cita duplicada o equivocada",      "Corrige en Salud Total", "o avisa admisión"),
        ("Paciente molesto o con queja",     "Transfiere a clínica",  "55 2884 1932"),
        ("Caso nuevo / aseguradora / pase",  "Deriva a admisión",     "55 5516 9900"),
        ("Dudas operativas internas",        "Dr. Ivan o Vicci",      "—"),
    ]
    y0 = Inches(2.1)
    add_card(s, Inches(0.6), y0, Inches(12.1), Inches(0.55), fill=GRAY_DARK)
    add_text(s, Inches(0.9), y0, Inches(5.5), Inches(0.55),
             "SITUACIÓN", size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(6.5), y0, Inches(3.5), Inches(0.55),
             "QUÉ HACER", size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(10.2), y0, Inches(2.4), Inches(0.55),
             "CONTACTO", size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    for i, (sit, action, contact) in enumerate(rows):
        y = y0 + Inches(0.6 + i*0.7)
        fill = GRAY_LIGHT if i % 2 == 0 else WHITE
        add_card(s, Inches(0.6), y, Inches(12.1), Inches(0.65), fill=fill)
        add_text(s, Inches(0.9), y, Inches(5.5), Inches(0.65),
                 sit, size=14, color=GRAY_DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(6.5), y, Inches(3.5), Inches(0.65),
                 action, size=14, bold=True, color=TEAL_DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(10.2), y, Inches(2.4), Inches(0.65),
                 contact, size=14, color=CORAL, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    # ----- Slide 13 — Preguntas
    s = prs.slides.add_slide(blank); set_background(s, WHITE)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    band.line.fill.background()
    band.fill.solid(); band.fill.fore_color.rgb = TEAL_DARK
    add_text(s, Inches(0.6), Inches(2.6), Inches(12.1), Inches(1.2),
             "¿Preguntas?", size=72, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(0.6),
             "Antes de terminar, vamos por las dudas de todos.",
             size=20, color=RGBColor(0xCF, 0xE7, 0xEA), align=PP_ALIGN.CENTER)
    add_rule(s, Inches(6.0), Inches(5.0), Inches(1.3), color=CORAL, height_pt=4)
    add_text(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(0.5),
             "Gracias por hacer posible TherapyCord.",
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    prs.save(path)
    print(f"✅ PPTX → {path}")


# ===================== PDF Handout =====================

def build_handout(path: Path) -> None:
    c = canvas.Canvas(str(path), pagesize=LETTER)
    W, H = LETTER

    teal = HexColor("#0E7C86")
    teal_dark = HexColor("#084C53")
    coral = HexColor("#E86A4C")
    gray = HexColor("#2B2B2B")
    gray_mid = HexColor("#555555")
    gray_light = HexColor("#F2F2F2")

    # --- Page 1 ---
    # Header band
    c.setFillColor(teal)
    c.rect(0, H - 1.1*inch, W, 1.1*inch, stroke=0, fill=1)
    c.setFillColor(HexColor("#FFFFFF"))
    c.setFont("Helvetica-Bold", 20)
    c.drawString(0.6*inch, H - 0.65*inch, "TherapyCord · Cheatsheet del terapeuta")
    c.setFont("Helvetica", 11)
    c.drawString(0.6*inch, H - 0.9*inch, "Capacitación 24 abr 2026  ·  Salud Total + Cordelia")

    # Section 1 — Números clave
    y = H - 1.5*inch
    c.setFillColor(teal_dark); c.setFont("Helvetica-Bold", 13)
    c.drawString(0.6*inch, y, "NÚMEROS QUE TIENES QUE RECORDAR")
    c.setStrokeColor(teal); c.setLineWidth(1.5)
    c.line(0.6*inch, y - 3, 3.2*inch, y - 3)

    # Two boxes side by side
    box_y = y - 1.75*inch
    box_h = 1.6*inch
    # Admisión box
    c.setFillColor(teal)
    c.roundRect(0.6*inch, box_y, 3.5*inch, box_h, 8, stroke=0, fill=1)
    c.setFillColor(HexColor("#FFFFFF"))
    c.setFont("Helvetica-Bold", 10); c.drawString(0.8*inch, box_y + box_h - 0.28*inch, "ADMISIÓN")
    c.setFont("Helvetica-Bold", 22); c.drawString(0.8*inch, box_y + box_h - 0.7*inch, "55 5516 9900")
    c.setFont("Helvetica", 9)
    for i, t in enumerate(["Pacientes nuevos", "Aseguradora / GMM",
                            "Pase médico", "Reembolso"]):
        c.drawString(0.8*inch, box_y + box_h - (0.95 + i*0.16)*inch, "• " + t)

    # Clínica box
    c.setFillColor(gray_light)
    c.roundRect(4.3*inch, box_y, 3.5*inch, box_h, 8, stroke=0, fill=1)
    c.setFillColor(coral)
    c.setFont("Helvetica-Bold", 10); c.drawString(4.5*inch, box_y + box_h - 0.28*inch, "CLÍNICA (CELULAR)")
    c.setFillColor(teal_dark)
    c.setFont("Helvetica-Bold", 22); c.drawString(4.5*inch, box_y + box_h - 0.7*inch, "55 2884 1932")
    c.setFillColor(gray)
    c.setFont("Helvetica", 9)
    for i, t in enumerate(["Queja o frustración", "Emergencia",
                            "Transferencia a humano", "Caso que Cordelia no resuelve"]):
        c.drawString(4.5*inch, box_y + box_h - (0.95 + i*0.16)*inch, "• " + t)

    # Section 2 — Regla de oro
    y = box_y - 0.4*inch
    c.setFillColor(teal_dark); c.setFont("Helvetica-Bold", 13)
    c.drawString(0.6*inch, y, "REGLA DE ORO")
    c.setStrokeColor(teal); c.setLineWidth(1.5)
    c.line(0.6*inch, y - 3, 2.1*inch, y - 3)
    c.setFillColor(gray); c.setFont("Helvetica", 11)
    c.drawString(0.6*inch, y - 0.35*inch,
                 "Cordelia SOLO agenda a pacientes EXISTENTES.")
    c.drawString(0.6*inch, y - 0.58*inch,
                 "Nuevos, aseguradora, pase o reembolso → ADMISIÓN 55 5516 9900.")

    # Section 3 — Cómo valida Cordelia
    y = y - 1.0*inch
    c.setFillColor(teal_dark); c.setFont("Helvetica-Bold", 13)
    c.drawString(0.6*inch, y, "CÓMO VALIDA CORDELIA A UN PACIENTE EXISTENTE")
    c.line(0.6*inch, y - 3, 3.8*inch, y - 3)

    c.setFillColor(gray); c.setFont("Helvetica", 10)
    c.drawString(0.6*inch, y - 0.3*inch,
                 "Cordelia pide DOS datos y los compara contra Salud Total:")
    c.setFont("Helvetica-Bold", 11); c.setFillColor(teal)
    c.drawString(0.7*inch, y - 0.55*inch, "1.  Teléfono (siempre)")
    c.setFillColor(gray_mid); c.setFont("Helvetica", 11)
    c.drawString(0.7*inch, y - 0.75*inch, "       MÁS")
    c.setFillColor(teal); c.setFont("Helvetica-Bold", 11)
    c.drawString(0.7*inch, y - 0.95*inch, "2.  Fecha de nacimiento   O   correo electrónico")
    c.setFillColor(coral); c.setFont("Helvetica-Oblique", 10)
    c.drawString(0.7*inch, y - 1.2*inch,
                 "Si no coinciden → se deriva a admisión. Nunca se agenda a ciegas.")

    # Section 4 — Tu checklist diario
    y = y - 1.7*inch
    c.setFillColor(teal_dark); c.setFont("Helvetica-Bold", 13)
    c.drawString(0.6*inch, y, "TU CHECKLIST DIARIO")
    c.line(0.6*inch, y - 3, 2.7*inch, y - 3)
    checklist = [
        "Abrir tu agenda en Salud Total al iniciar el turno",
        "Revisar notas de sesión previa de cada paciente",
        "Atender sesión 40-45 min individualizados",
        "Registrar nota de sesión al terminar (obligatorio)",
        "Avisar a recepción si hay cambios de horario",
    ]
    c.setFillColor(gray); c.setFont("Helvetica", 10)
    for i, item in enumerate(checklist):
        # checkbox
        c.setStrokeColor(teal); c.setLineWidth(1)
        c.rect(0.65*inch, y - 0.4*inch - i*0.25*inch - 2,
               0.15*inch, 0.15*inch, stroke=1, fill=0)
        c.drawString(0.92*inch, y - 0.4*inch - i*0.25*inch, item)

    # Footer page 1
    c.setFillColor(gray_mid); c.setFont("Helvetica", 8)
    c.drawCentredString(W/2, 0.4*inch,
                         "TherapyCord · Hospital Ángeles México · Torre B, Piso 7, Consultorio 751")
    c.drawCentredString(W/2, 0.25*inch, "Página 1 de 2")

    c.showPage()

    # --- Page 2 ---
    # Header
    c.setFillColor(teal)
    c.rect(0, H - 1.1*inch, W, 1.1*inch, stroke=0, fill=1)
    c.setFillColor(HexColor("#FFFFFF"))
    c.setFont("Helvetica-Bold", 20)
    c.drawString(0.6*inch, H - 0.65*inch, "Accesos y soporte")
    c.setFont("Helvetica", 11)
    c.drawString(0.6*inch, H - 0.9*inch, "Esta hoja es tuya. Guárdala cerca del consultorio.")

    # Salud Total
    y = H - 1.5*inch
    c.setFillColor(teal_dark); c.setFont("Helvetica-Bold", 13)
    c.drawString(0.6*inch, y, "SALUD TOTAL")
    c.setStrokeColor(teal); c.setLineWidth(1.5)
    c.line(0.6*inch, y - 3, 1.8*inch, y - 3)

    c.setFillColor(gray); c.setFont("Helvetica", 11)
    c.drawString(0.6*inch, y - 0.3*inch, "URL:  www.saludtotal.mx")
    c.drawString(0.6*inch, y - 0.52*inch,
                 "Contraseña inicial: 123456789  (cámbiala en tu primer acceso)")

    # Table of accesses
    rows = [
        ("Lic. Lilia Salazar",    "lilia@therapycord.com"),
        ("Lic. Montserrat",       "monste@therapycord.com"),
        ("Lic. Kevin Castellanos","kevin@therapycord.com"),
        ("Lic. Harold Ildefonso", "harold@therapycord.com"),
    ]
    ty = y - 0.95*inch
    c.setFillColor(HexColor("#2B2B2B"))
    c.rect(0.6*inch, ty, 6.8*inch, 0.3*inch, stroke=0, fill=1)
    c.setFillColor(HexColor("#FFFFFF"))
    c.setFont("Helvetica-Bold", 10)
    c.drawString(0.8*inch, ty + 0.1*inch, "TERAPEUTA")
    c.drawString(3.4*inch, ty + 0.1*inch, "USUARIO")
    for i, (name, email) in enumerate(rows):
        row_y = ty - (i+1)*0.3*inch
        c.setFillColor(gray_light if i % 2 == 0 else HexColor("#FFFFFF"))
        c.rect(0.6*inch, row_y, 6.8*inch, 0.3*inch, stroke=0, fill=1)
        c.setFillColor(gray); c.setFont("Helvetica-Bold", 10)
        c.drawString(0.8*inch, row_y + 0.1*inch, name)
        c.setFont("Helvetica", 10)
        c.drawString(3.4*inch, row_y + 0.1*inch, email)

    # Equipo y horarios
    y = ty - 0.3*inch*5 - 0.3*inch
    c.setFillColor(teal_dark); c.setFont("Helvetica-Bold", 13)
    c.drawString(0.6*inch, y, "QUIÉN ATIENDE QUÉ")
    c.line(0.6*inch, y - 3, 2.3*inch, y - 3)

    matrix = [
        ("🌅 TURNO MAÑANA  (8am-3pm)",
         [("Lic. Lilia", "Suelo pélvico, pre/post parto"),
          ("Lic. Montserrat", "Ortopédica general")]),
        ("🌆 TURNO TARDE  (1pm-8pm)",
         [("Lic. Kevin", "Deportiva, neurológica"),
          ("Lic. Harold", "Ortopédica, post-COVID")]),
    ]
    ty = y - 0.3*inch
    for header, rows in matrix:
        c.setFillColor(teal); c.setFont("Helvetica-Bold", 11)
        c.drawString(0.6*inch, ty, header)
        ty -= 0.22*inch
        c.setFillColor(gray); c.setFont("Helvetica", 10)
        for name, spec in rows:
            c.setFont("Helvetica-Bold", 10)
            c.drawString(0.8*inch, ty, name + " —")
            c.setFont("Helvetica", 10)
            c.drawString(2.0*inch, ty, spec)
            ty -= 0.2*inch
        ty -= 0.1*inch

    # Si algo falla
    y = ty - 0.1*inch
    c.setFillColor(teal_dark); c.setFont("Helvetica-Bold", 13)
    c.drawString(0.6*inch, y, "SI ALGO FALLA")
    c.line(0.6*inch, y - 3, 2.1*inch, y - 3)

    trouble = [
        ("Cordelia no contestó una llamada",   "Vicci · 55 1188 0301"),
        ("Cita duplicada o equivocada",        "Admisión · 55 5516 9900"),
        ("Paciente molesto o con queja",       "Clínica · 55 2884 1932"),
        ("Caso nuevo / aseguradora / pase",    "Admisión · 55 5516 9900"),
        ("Dudas operativas internas",          "Dr. Ivan o Vicci"),
    ]
    ty = y - 0.3*inch
    c.setFillColor(HexColor("#2B2B2B"))
    c.rect(0.6*inch, ty, 7.3*inch, 0.3*inch, stroke=0, fill=1)
    c.setFillColor(HexColor("#FFFFFF")); c.setFont("Helvetica-Bold", 10)
    c.drawString(0.8*inch, ty + 0.1*inch, "SITUACIÓN")
    c.drawString(4.5*inch, ty + 0.1*inch, "CONTACTO")
    for i, (sit, cnt) in enumerate(trouble):
        row_y = ty - (i+1)*0.3*inch
        c.setFillColor(gray_light if i % 2 == 0 else HexColor("#FFFFFF"))
        c.rect(0.6*inch, row_y, 7.3*inch, 0.3*inch, stroke=0, fill=1)
        c.setFillColor(gray); c.setFont("Helvetica", 10)
        c.drawString(0.8*inch, row_y + 0.1*inch, sit)
        c.setFillColor(coral); c.setFont("Helvetica-Bold", 10)
        c.drawString(4.5*inch, row_y + 0.1*inch, cnt)

    # Footer page 2
    c.setFillColor(gray_mid); c.setFont("Helvetica", 8)
    c.drawCentredString(W/2, 0.4*inch,
                         "TherapyCord · v1 · 24 abril 2026")
    c.drawCentredString(W/2, 0.25*inch, "Página 2 de 2")

    c.save()
    print(f"✅ PDF  → {path}")


if __name__ == "__main__":
    build_pptx(OUT_DIR / "TherapyCord_Capacitacion_Terapeutas.pptx")
    build_handout(OUT_DIR / "TherapyCord_Cheatsheet_Terapeutas.pdf")
